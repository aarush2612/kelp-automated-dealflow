from pptx import Presentation


FOOTER_TEXT = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
LOGO_TEXT = "Kelp"


def test_footer_and_logo_present():
    prs = Presentation("outputs/entertainment-connplex/TargetCo_Teaser.pptx")

    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

        full_text = " ".join(texts)

        assert FOOTER_TEXT in full_text, "Footer missing on slide"
        assert LOGO_TEXT in full_text, "Logo text missing on slide"
