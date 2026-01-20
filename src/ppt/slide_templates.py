from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_footer(slide):
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(6.8),
        Inches(8), Inches(0.3)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
    p.font.size = Pt(9)
    p.alignment = PP_ALIGN.CENTER


def add_title(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(28)


def add_bullets(slide, text_list, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()

    for i, txt in enumerate(text_list):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(14)
