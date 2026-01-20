from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


class ChartBuilder:

    @staticmethod
    def revenue_chart(slide, revenue_data: dict):
        chart_data = ChartData()
        years = list(revenue_data.keys())
        values = [revenue_data[y] for y in years]

        chart_data.categories = years
        chart_data.add_series("Revenue (₹ Cr)", values)

        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(4.5), Inches(3)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, cx, cy,
            chart_data
        )

    @staticmethod
    def pat_chart(slide, pat_data: dict):
        chart_data = ChartData()
        years = list(pat_data.keys())
        values = [pat_data[y] for y in years]

        chart_data.categories = years
        chart_data.add_series("PAT (₹ Cr)", values)

        x, y, cx, cy = Inches(5.2), Inches(1.5), Inches(4.3), Inches(3)
        slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            x, y, cx, cy,
            chart_data
        )
