from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

from src.ppt.slide_templates import add_footer, add_title, add_bullets
from src.ppt.chart_builder import ChartBuilder


class PPTBuilder:

    def __init__(self):
        self.prs = Presentation()

    def build(self, company_data: dict, insights: dict, output_path: str):
        self._slide_overview(company_data, insights)
        self._slide_financials(company_data)
        self._slide_swot(company_data)

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)

    # ---------------- slides ----------------

    def _slide_overview(self, data, insights):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        add_title(slide, "Company Overview")

        bullets = [
            data["company_profile"]["description"],
            f"Website: {data['company_profile']['website']}"
        ] + insights["highlights"]

        add_bullets(
            slide,
            bullets,
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(9),
            height=Inches(4)
        )

        add_footer(slide)

    def _slide_financials(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        add_title(slide, "Financial Performance")

        income = data["financials"].get("Income Statement", {})
        revenue = income.get("Revenue From Operations", {})
        pat = income.get("PAT", {})

        if revenue:
            ChartBuilder.revenue_chart(slide, revenue)
        if pat:
            ChartBuilder.pat_chart(slide, pat)

        add_footer(slide)

    def _slide_swot(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        add_title(slide, "SWOT Summary")

        swot = data.get("swot", {})
        bullets = []

        for k in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
            if k in swot:
                bullets.append(f"{k}: {', '.join(swot[k][:2])}")

        add_bullets(
            slide,
            bullets,
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(9),
            height=Inches(4)
        )

        add_footer(slide)
